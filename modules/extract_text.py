# modules/extract_text.py
import fitz  # PyMuPDF
from pptx import Presentation
import io
import tempfile
def extract_from_pdf(uploaded_file):
    """Extract all text from a PDF file uploaded via Streamlit."""
    text = ""
    # Use BytesIO object directly
    with fitz.open(stream=uploaded_file.read(), filetype="pdf") as pdf:
        for page in pdf:
            text += page.get_text("text")
    return text.strip()

def extract_from_pptx(uploaded_file):
    """Extract all text from a PPTX presentation."""
    text = ""
    presentation = Presentation(io.BytesIO(uploaded_file.read()))
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

def extract_from_text(uploaded_file):
    """Extract text from a plain .txt file."""
    return uploaded_file.read().decode("utf-8").strip()
import tempfile

def extract_from_pdf(uploaded_file):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(uploaded_file.read())
        tmp.flush()
        with fitz.open(tmp.name) as pdf:
            text = "".join(page.get_text() for page in pdf)
    return text.strip()
